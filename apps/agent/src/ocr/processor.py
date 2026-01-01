"""Document processor for extracting text from various file types."""

import logging
from pathlib import Path
from typing import List, Optional

from .paddle_ocr import PaddleOCRProcessor

logger = logging.getLogger("doc-processor")


class DocumentProcessor:
    """
    Unified document processor for extracting text from various formats.

    Supports:
    - PDF (text-based and scanned)
    - Images (PNG, JPG, etc.)
    - Word documents (DOCX)
    - PowerPoint (PPTX)
    - Plain text/Markdown
    """

    def __init__(self, use_gpu: bool = False, lang: str = "multilingual"):
        """
        Initialize document processor.

        Args:
            use_gpu: Whether to use GPU for OCR
            lang: OCR language setting
        """
        self.ocr = PaddleOCRProcessor(use_gpu=use_gpu, lang=lang)

    def extract_text(self, file_path: str, mime_type: Optional[str] = None) -> str:
        """
        Extract text from a file based on its type.

        Args:
            file_path: Path to the file
            mime_type: Optional MIME type hint

        Returns:
            Extracted text content
        """
        path = Path(file_path)
        suffix = path.suffix.lower()

        # Determine file type
        if mime_type:
            if "pdf" in mime_type:
                return self._process_pdf(file_path)
            elif "image" in mime_type:
                return self._process_image(file_path)
            elif "word" in mime_type or "docx" in mime_type:
                return self._process_docx(file_path)
            elif "powerpoint" in mime_type or "pptx" in mime_type:
                return self._process_pptx(file_path)
            elif "text" in mime_type or "markdown" in mime_type:
                return self._process_text(file_path)

        # Fall back to extension
        if suffix == ".pdf":
            return self._process_pdf(file_path)
        elif suffix in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"]:
            return self._process_image(file_path)
        elif suffix == ".docx":
            return self._process_docx(file_path)
        elif suffix == ".pptx":
            return self._process_pptx(file_path)
        elif suffix in [".txt", ".md", ".markdown"]:
            return self._process_text(file_path)
        else:
            logger.warning(f"Unknown file type: {suffix}")
            return self._process_text(file_path)  # Try as text

    def _process_pdf(self, file_path: str) -> str:
        """Extract text from PDF, using OCR if needed."""
        try:
            import fitz

            doc = fitz.open(file_path)
            text_parts = []

            for page in doc:
                page_text = page.get_text()
                text_parts.append(page_text)

            doc.close()
            text = "\n\n".join(text_parts)

            # If very little text, it might be scanned - use OCR
            if len(text.strip()) < 100 and self.ocr.needs_ocr(file_path):
                logger.info(f"PDF appears scanned, using OCR: {file_path}")
                return self.ocr.ocr_pdf_to_text(file_path)

            return text

        except Exception as e:
            logger.error(f"PDF processing error: {e}")
            # Fall back to OCR
            return self.ocr.ocr_pdf_to_text(file_path)

    def _process_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        text, confidence = self.ocr.ocr_image(file_path)
        logger.info(f"Image OCR confidence: {confidence:.2f}")
        return text

    def _process_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)

        except ImportError:
            logger.error("python-docx not installed. Run: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"DOCX processing error: {e}")
            return ""

    def _process_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)

        except ImportError:
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"PPTX processing error: {e}")
            return ""

    def _process_text(self, file_path: str) -> str:
        """Read plain text file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, "r", encoding="latin-1") as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Text file read error: {e}")
                return ""
        except Exception as e:
            logger.error(f"Text file read error: {e}")
            return ""


def chunk_text(
    text: str,
    chunk_size: int = 500,
    overlap: int = 50,
) -> List[str]:
    """
    Split text into overlapping chunks for indexing.

    Args:
        text: The text to chunk
        chunk_size: Target size in tokens (approximately chars/4)
        overlap: Number of tokens to overlap between chunks

    Returns:
        List of text chunks
    """
    # Approximate tokens as chars/4
    char_chunk_size = chunk_size * 4
    char_overlap = overlap * 4

    if len(text) <= char_chunk_size:
        return [text]

    chunks = []
    start = 0

    while start < len(text):
        end = start + char_chunk_size

        # Try to break at sentence boundary
        if end < len(text):
            # Look for sentence end within last 20% of chunk
            search_start = end - int(char_chunk_size * 0.2)
            search_region = text[search_start:end]

            for sep in [". ", "! ", "? ", "\n\n", "\n"]:
                last_sep = search_region.rfind(sep)
                if last_sep != -1:
                    end = search_start + last_sep + len(sep)
                    break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - char_overlap

    return chunks
